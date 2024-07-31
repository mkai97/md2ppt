import os
from fastapi import APIRouter
import uvicorn
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
import random
import string

router = APIRouter()

@router.post("/convert")
def create_ppt(presentation: dict,pptInfo:dict):

    # 合并数据
   # 检查presentation是否为空，如果为空，则返回一个包含提示信息的字典
    if presentation == None:
        return {"message": "数据为空"}

    # 如果presentation不为空，开始处理数据
    else:
        # 将presentation赋值给data，以便后续操作
        data = presentation
        
        # 获取幻灯片列表
        slides = data["slides"]
        
        # 更新pptInfo字典，添加副标题、作者和创建日期信息
        pptInfo["subtitle"] = pptInfo.get("subtitle","")+"\n 作者："+pptInfo.get("author","")+"\n 日期："+pptInfo.get("createTime","")
        
        # 初始化幻灯片的布局和内容
        pptInfo["layout"] = 0
        pptInfo["content"] = ""
        
    # 将pptInfo作为新的幻灯片插入到幻灯片列表的开头
    slides.insert(0,pptInfo)

        # 结束语
        
    slides.append( {
                "layout": 1,
                "title": "感谢！",
                "subtitle": "感谢您的关注和支持，今天的演讲到此结束，谢谢大家。"
                ,
                "content":[],
                "image": {
                    "path": "bg01.jpg",
                    "size": {
                        "width": 10,
                        "height": 7.5
                    }
                }
            })





    
    # 更新data字典中的幻灯片列表
    data["slides"] = slides
 
    # 创建一个新的 PPTX 演示文稿
    prs = Presentation()


    # 遍历 JSON 数据中的幻灯片
    for slide_data in data["slides"]:
        slide_layout = prs.slide_layouts[slide_data["layout"]]
        slide = prs.slides.add_slide(slide_layout)

        # 设置标题和子标题
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        tf = subtitle.text_frame
        sp = tf.add_paragraph()
        title.text = slide_data["title"]
        sp.text = slide_data.get("subtitle", "")

        sp.font.color.rgb = RGBColor(255, 255, 255)
        sp.font.size = Pt(24)


        # 设置标题和子标题的字体颜色为白色
        p = title.text_frame.paragraphs[0]
        run = p.runs[0]
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.size = Pt(34)



        if "image" in data:
            current_file_path = os.path.realpath(__file__)
            current_dir = os.path.dirname(current_file_path)
            img_path = os.path.join(current_dir, '..', 'pptSource', 'static', data["image"]["path"])

            # img_path = data["image"]["path"]
            width = Inches(data["image"]["size"]["width"])
            height = Inches(data["image"]["size"]["height"])
            pic = slide.shapes.add_picture(img_path, 0, 0, width=width, height=height)
            pic.line.fill.background()

            # 将图片放到所有形状的最底层
            spTree = slide.shapes._spTree
            spTree.remove(pic._element)
            spTree.insert(0, pic._element)


        # 添加图片
        if "image" in slide_data:
            current_file_path = os.path.realpath(__file__)
            current_dir = os.path.dirname(current_file_path)
            img_path = os.path.join(current_dir, '..', 'pptSource', 'static',  slide_data["image"]["path"])
            width = Inches(slide_data["image"]["size"]["width"])
            height = Inches(slide_data["image"]["size"]["height"])
            pic = slide.shapes.add_picture(img_path, 0, 0, width=width, height=height)
            pic.line.fill.background()

            # 将图片放到所有形状的最底层
            spTree = slide.shapes._spTree
            spTree.remove(pic._element)
            spTree.insert(1, pic._element)

        # 添加内容文本
        tf = slide.placeholders[1].text_frame
        for content in slide_data["content"]:
            p = tf.add_paragraph()
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.size = Pt(14)
            p.text = content
            p.level = 1

    # 保存演示文稿
    current_file_path = os.path.realpath(__file__)
    current_dir = os.path.dirname(current_file_path)
    file_name = presentation["title"]+"-"+generate_random_string()+".pptx"
    ppt_path = os.path.join(current_dir, '..', 'pptSource', 'ppt', file_name)
    prs.save(ppt_path)
    return {"code":0,
        "message": "success",
            "file_path": service_root() +"/pptSource/ppt/"+ file_name}

# 定义一个函数来生成 6 位随机字符的字符串
def generate_random_string(length=6):
    # 随机选择 ASCII 字母和数字
    characters = string.ascii_letters + string.digits
    # 使用 random.choice 从字符中随机选择一个字符，并重复 length 次
    random_string = ''.join(random.choice(characters) for _ in range(length))
    return random_string

def service_root():
    host = "192.168.1.111"  # 默认是 '127.0.0.1'
    port = "8000"  # 默认是 8000
    return  "http://"+ f"{host}:{port}"

