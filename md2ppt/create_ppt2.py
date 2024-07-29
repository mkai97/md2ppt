import json

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches


def create_ppt(context):
    # 加载 JSON 数据
    with open('ppt2.json', 'r') as file:
        data = json.load(file)

    # 创建一个新的 PPTX 演示文稿
    prs = Presentation()

    # 设置演示
    # 加载 JSON 数据
    with open('presentation.json', 'r') as file:
        data = json.load(file)

    # 创建一个新的 PPTX 演示文稿
    prs = Presentation()

    # 遍历 JSON 数据中的幻灯片
    for slide_data in data["presentation"]["slides"]:
        slide_layout = prs.slide_layouts[slide_data["layout"]]
        slide = prs.slides.add_slide(slide_layout)

        # 设置标题和子标题
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = slide_data["title"]
        subtitle.text = slide_data.get("subtitle", "")

        # 设置标题和子标题的字体颜色为白色
        p = title.text_frame.paragraphs[0]
        run = p.runs[0]
        run.font.color.rgb = RGBColor(255, 255, 255)

        # 添加图片
        if "image" in slide_data:
            img_path = slide_data["image"]["path"]
            width = Inches(slide_data["image"]["size"]["width"])
            height = Inches(slide_data["image"]["size"]["height"])
            pic = slide.shapes.add_picture(img_path, 0, 0, width=width, height=height)
            pic.line.fill.background()

            # 将图片放到所有形状的最底层
            spTree = slide.shapes._spTree
            spTree.remove(pic._element)
            spTree.insert(0, pic._element)

        # 添加内容文本
        tf = slide.placeholders[2].text_frame
        for content in slide_data["content"]:
            p = tf.add_paragraph()
            p.text = content
            p.level = 1

    # 保存演示文稿
    prs.save('雷军演讲2024.pptx')
