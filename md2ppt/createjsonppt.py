import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# 定义一个函数来设置所有文本的颜色为白色
def set_text_color_white(presentation):
    for slide in presentation.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)

# 加载JSON数据
json_data = '''
[
  {
    "slide_number": 1,
    "title": "2024年雷军年度演讲",
    "content": [
      "欢迎词",
      "演讲主题介绍",
      "2024年科技趋势预测"
    ]
  },
  {
    "slide_number": 2,
    "title": "小米的2023年回顾",
    "content": [
      "2023年小米的主要成就",
      "产品创新亮点",
      "市场表现与用户反馈"
    ]
  },
  {
    "slide_number": 3,
    "title": "2024年小米战略规划",
    "content": [
      "新产品线发布计划",
      "技术革新与研发方向",
      "市场拓展与品牌战略"
    ]
  },
  {
    "slide_number": 4,
    "title": "AI与IoT的融合",
    "content": [
      "AI技术在小米产品中的应用",
      "IoT生态系统的构建与优化",
      "未来智能家居的展望"
    ]
  },
  {
    "slide_number": 5,
    "title": "社会责任与可持续发展",
    "content": [
      "小米的社会责任实践",
      "可持续发展战略",
      "环保与公益项目介绍"
    ]
  }
]
'''

# 解析JSON数据
data = json.loads(json_data)

# 初始化演示文稿
prs = Presentation()

# 循环处理每一项数据
for item in data:
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = item['title']
    
    body = slide.placeholders[1]
    tf = body.text_frame
    
    for content_item in item['content']:
        p = tf.add_paragraph()
        p.text = content_item
        p.level = 1

# 应用背景图片 (假设有一个背景图片文件)
img_path = 'background_image.png'  # 背景图片路径

# 添加背景图片
for slide in prs.slides:
    pic = slide.shapes.add_picture(img_path, 0, 0, width=Inches(10), height=Inches(7.5))
    pic.line.fill.background()
    # 将图片放到所有形状的最底层
    spTree = slide.shapes._spTree
    spTree.remove(pic._element)
    spTree.insert(0, pic._element)

# 调整字体大小（可选）
for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                font = run.font
                font.size = Pt(18)

# 调用函数设置所有文本颜色为白色
set_text_color_white(prs)

# 保存演示文稿
prs.save('dynamic_presentation.pptx')