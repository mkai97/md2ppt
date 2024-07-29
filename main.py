from fastapi import FastAPI
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

app = FastAPI()


@app.get("/")
def read_root():
    return {"Hello": "World"}


@app.post("/md2ppt")
def create_ppt(presentation: dict,pptInfo:dict):
   # 加载 JSON 数据
    # with open('./md2ppt/ppt2.json', 'r') as file:
        # data = json.load(file)

    # 合并数据

   # 检查presentation是否为空，如果为空，则返回一个包含提示信息的字典
    if presentation == None:
        return {"message": "数据为空"}

    # 如果presentation不为空，开始处理数据
    else:
        # 将presentation赋值给data，以便后续操作
        data = presentation
        
        # 获取幻灯片列表
        slides = data["slides"]
        
        # 更新pptInfo字典，添加副标题、作者和创建日期信息
        pptInfo["subtitle"] = pptInfo.get("subtitle","")+"\n 作者："+pptInfo.get("author","")+"\n 日期："+pptInfo.get("createTime","")
        
        # 初始化幻灯片的布局和内容
        pptInfo["layout"] = 0
        pptInfo["content"] = ""
        
    # 将pptInfo作为新的幻灯片插入到幻灯片列表的开头
    slides.insert(0,pptInfo)
    
    # 更新data字典中的幻灯片列表
    data["slides"] = slides
 
    # 创建一个新的 PPTX 演示文稿
    prs = Presentation()


    # 遍历 JSON 数据中的幻灯片
    for slide_data in data["slides"]:
        slide_layout = prs.slide_layouts[slide_data["layout"]]
        slide = prs.slides.add_slide(slide_layout)

        # 设置标题和子标题
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        tf = subtitle.text_frame
        sp = tf.add_paragraph()
        title.text = slide_data["title"]
        sp.text = slide_data.get("subtitle", "")

        sp.font.color.rgb = RGBColor(255, 255, 255)
        sp.font.size = Pt(24)


        # 设置标题和子标题的字体颜色为白色
        p = title.text_frame.paragraphs[0]
        run = p.runs[0]
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.size = Pt(34)



        if "image" in data:
            img_path = data["image"]["path"]
            width = Inches(data["image"]["size"]["width"])
            height = Inches(data["image"]["size"]["height"])
            pic = slide.shapes.add_picture(img_path, 0, 0, width=width, height=height)
            pic.line.fill.background()

            # 将图片放到所有形状的最底层
            spTree = slide.shapes._spTree
            spTree.remove(pic._element)
            spTree.insert(0, pic._element)


        # 添加图片
        if "image" in slide_data:
            img_path = slide_data["image"]["path"]
            width = Inches(slide_data["image"]["size"]["width"])
            height = Inches(slide_data["image"]["size"]["height"])
            pic = slide.shapes.add_picture(img_path, 0, 0, width=width, height=height)
            pic.line.fill.background()

            # 将图片放到所有形状的最底层
            spTree = slide.shapes._spTree
            spTree.remove(pic._element)
            spTree.insert(1, pic._element)

        # 添加内容文本
        tf = slide.placeholders[1].text_frame
        for content in slide_data["content"]:
            p = tf.add_paragraph()
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.size = Pt(14)
            p.text = content
            p.level = 1

    # 保存演示文稿
    prs.save('雷军演讲2024.pptx')


@app.get("/items/{item_id}")
def read_item(item_id: int, q: str = None):
    return {"item_id": item_id, "q": q}
