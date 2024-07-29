from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


def set_text_color_white(presentation):
    """递归遍历演示文稿的所有文本框并设置文本颜色为白色。"""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)


# 初始化演示文稿
prs = Presentation()

# 应用背景图片
img_path = 'background_image.png'  # 背景图片路径

# 幻灯片1: 封面
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "2024年度雷军演讲"
subtitle.text = "探索科技与创新的未来\n演讲者：雷军\n演讲日期：2024年7月22日\n地点：[具体地点]"
p = title.text_frame.paragraphs[0]
run = p.runs[0]
run.font.color.rgb = RGBColor(255, 255, 255)  # 设置字体颜色为白色

# 添加图片
pic = slide.shapes.add_picture(img_path, 0, 0, width=Inches(10), height=Inches(7.5))
pic.line.fill.background()

# 将图片放到所有形状的最底层
spTree = slide.shapes._spTree
spTree.remove(pic._element)
spTree.insert(0, pic._element)

# 幻灯片2: 欢迎与介绍
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "勇气从何而来"
body = slide.placeholders[1]
tf = body.text_frame
tf.text = "尊敬的听众，大家好！我是雷军，非常荣幸能在这里与大家分享我的一些想法和见解。"
p = tf.add_paragraph()
p.text = "小米科技创始人，致力于推动科技与创新的融合。"
p.level = 1
p = tf.add_paragraph()
p.text = "今天，我将与大家探讨科技行业的发展趋势，以及小米科技在这一过程中的角色和贡献。"
p.level = 1

# 幻灯片3: 2024年科技行业回顾
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
body = slide.placeholders[1]
title.text = "2024年科技行业回顾"
tf = body.text_frame
tf.text = "今年，科技行业经历了前所未有的变革，包括[具体事件]。"
p = tf.add_paragraph()
p.text = "我们看到[具体趋势]正在塑造我们的未来。"
p.level = 1

# 幻灯片4: 小米科技的2024年成就
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
body = slide.placeholders[1]
title.text = "小米科技的2024年成就"
tf = body.text_frame
tf.text = "产品创新：今年，我们推出了[具体产品]，它代表了我们对创新的不懈追求。"
p = tf.add_paragraph()
p.text = "市场扩张：我们的市场已经扩展到[具体地区]，我们的产品受到了广泛欢迎。"
p.level = 1
p = tf.add_paragraph()
p.text = "技术突破：在[具体技术]领域，我们取得了突破性进展。"
p.level = 1

# 幻灯片5: 雷军的个人故事与哲学
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
body = slide.placeholders[1]
title.text = "雷军的个人故事与哲学"
tf = body.text_frame
tf.text = "从[具体经历]到今天，我一直坚信[具体哲学]。"
p = tf.add_paragraph()
p.text = "我相信[具体价值观]是推动我们前进的关键。"
p.level = 1

# 幻灯片6: 面对挑战与机遇
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
body = slide.placeholders[1]
title.text = "面对挑战与机遇"
tf = body.text_frame
tf.text = "当前，我们面临[具体挑战]。"
p = tf.add_paragraph()
p.text = "小米科技通过[具体策略]来应对这些挑战。"
p.level = 1

# 幻灯片7: 未来展望
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
body = slide.placeholders[1]
title.text = "未来展望"
tf = body.text_frame
tf.text = "我们计划[具体计划]，以实现我们的目标。"
p = tf.add_paragraph()
p.text = "我预测，[具体预测]将是未来科技行业的关键趋势。"
p.level = 1

# 幻灯片8: 教育与创新
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
body = slide.placeholders[1]
title.text = "教育与创新"
tf = body.text_frame
tf.text = "教育是培养创新精神的摇篮。"
p = tf.add_paragraph()
p.text = "小米科技一直致力于[具体贡献]，以支持教育和创新。"
p.level = 1

# 幻灯片9: 互动环节
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
body = slide.placeholders[1]
title.text = "互动环节"
tf = body.text_frame
tf.text = "现在，我非常期待与大家的互动。请随时提问。"
p = tf.add_paragraph()
p.text = "对于在座的年轻学生，我的建议是[具体建议]。"
p.level = 1

# 幻灯片10: 结束语和致谢
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
body = slide.placeholders[1]
title.text = "结束语和致谢"
tf = body.text_frame
tf.text = "感谢大家的聆听，希望今天的分享能给大家带来启发。"
p = tf.add_paragraph()
p.text = "再次感谢各位的参与，期待与大家在未来有更多的交流。"
p.level = 1

for slide in prs.slides:
    pic = slide.shapes.add_picture(img_path, 0, 0, width=Inches(10), height=Inches(7.5))
    pic.line.fill.background()
    # 将图片放到所有形状的最底层
    spTree = slide.shapes._spTree
    spTree.remove(pic._element)
    spTree.insert(0, pic._element)

# 调整字体大小（可选）
for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                font = run.font
                font.size = Pt(18)

# 调用函数设置所有文本颜色为白色
set_text_color_white(prs)

# 保存演示文稿
prs.save('lei_jun_speech_9.pptx')
